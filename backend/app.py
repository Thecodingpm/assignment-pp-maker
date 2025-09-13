from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
import base64
import tempfile
import os
import zipfile
import xml.etree.ElementTree as ET
import json
from io import BytesIO
from PIL import Image

app = Flask(__name__)
CORS(app)

def compress_image(image_data, max_size_kb=100, quality=85):
    """
    Compress image data to reduce size for Firebase storage
    """
    try:
        # Open image from bytes
        image = Image.open(BytesIO(image_data))
        
        # Convert to RGB if necessary (for JPEG)
        if image.mode in ('RGBA', 'LA', 'P'):
            # Create white background for transparent images
            background = Image.new('RGB', image.size, (255, 255, 255))
            if image.mode == 'P':
                image = image.convert('RGBA')
            background.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
            image = background
        elif image.mode != 'RGB':
            image = image.convert('RGB')
        
        # Calculate target size
        original_size_kb = len(image_data) / 1024
        if original_size_kb <= max_size_kb:
            return image_data, original_size_kb
        
        # Compress with quality adjustment
        output = BytesIO()
        current_quality = quality
        
        while current_quality > 10:
            output.seek(0)
            output.truncate(0)
            image.save(output, format='JPEG', quality=current_quality, optimize=True)
            compressed_size_kb = len(output.getvalue()) / 1024
            
            if compressed_size_kb <= max_size_kb:
                print(f"✅ Compressed image: {original_size_kb:.1f}KB → {compressed_size_kb:.1f}KB (quality: {current_quality})")
                return output.getvalue(), compressed_size_kb
            
            current_quality -= 10
        
        # If still too large, resize the image
        if compressed_size_kb > max_size_kb:
            # Calculate resize factor
            resize_factor = (max_size_kb / compressed_size_kb) ** 0.5
            new_size = (int(image.width * resize_factor), int(image.height * resize_factor))
            image = image.resize(new_size, Image.Resampling.LANCZOS)
            
            output.seek(0)
            output.truncate(0)
            image.save(output, format='JPEG', quality=70, optimize=True)
            final_size_kb = len(output.getvalue()) / 1024
            print(f"✅ Resized and compressed image: {original_size_kb:.1f}KB → {final_size_kb:.1f}KB")
            return output.getvalue(), final_size_kb
        
        return output.getvalue(), compressed_size_kb
        
    except Exception as e:
        print(f"❌ Error compressing image: {e}")
        return image_data, len(image_data) / 1024

def optimize_content_for_firebase(content, max_size_mb=0.9):
    """
    Optimize content size to fit within Firebase's 1MB limit
    """
    try:
        # Convert to JSON to get size
        json_str = json.dumps(content)
        current_size_mb = len(json_str) / (1024 * 1024)
        
        print(f"📊 Content size: {current_size_mb:.2f} MB")
        
        if current_size_mb <= max_size_mb:
            return content, current_size_mb
        
        print(f"⚠️ Content too large ({current_size_mb:.2f} MB), optimizing...")
        
        # If too large, try to reduce image quality further
        if 'slides' in content:
            for slide in content['slides']:
                if 'elements' in slide:
                    for element in slide['elements']:
                        if element.get('type') == 'image' and 'src' in element:
                            # Extract base64 data
                            src = element['src']
                            if src.startswith('data:image/jpeg;base64,'):
                                base64_data = src.split(',')[1]
                                try:
                                    # Decode, compress more aggressively, re-encode
                                    image_data = base64.b64decode(base64_data)
                                    compressed_data, _ = compress_image(image_data, max_size_kb=50, quality=60)
                                    new_base64 = base64.b64encode(compressed_data).decode('utf-8')
                                    element['src'] = f"data:image/jpeg;base64,{new_base64}"
                                    print(f"🔄 Further compressed image: {len(image_data)/1024:.1f}KB → {len(compressed_data)/1024:.1f}KB")
                                except Exception as e:
                                    print(f"❌ Error re-compressing image: {e}")
        
        # Check final size
        final_json = json.dumps(content)
        final_size_mb = len(final_json) / (1024 * 1024)
        print(f"✅ Optimized size: {final_size_mb:.2f} MB")
        
        return content, final_size_mb
        
    except Exception as e:
        print(f"❌ Error optimizing content: {e}")
        return content, len(json.dumps(content)) / (1024 * 1024)

def map_zip_images_to_slides(pptx_path, images):
    """
    Map ZIP-extracted images to their correct slide positions by analyzing slide XML
    """
    try:
        print(f"🔍 Mapping {len(images)} images to slide positions...")
        
        with zipfile.ZipFile(pptx_path, 'r') as zip_file:
            # Get all slide files
            slide_files = [f for f in zip_file.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
            print(f"📄 Found {len(slide_files)} slide XML files")
            
            # Create a mapping of media files to their usage in slides
            media_to_slides = {}
            
            for slide_file in slide_files:
                try:
                    # Extract slide number
                    if 'slides/' in slide_file:
                        slide_num = int(slide_file.split('slides/slide')[1].split('.')[0]) - 1
                    else:
                        slide_num = int(slide_file.split('slide')[1].split('.')[0]) - 1
                    
                    # Read slide XML
                    slide_xml = zip_file.read(slide_file)
                    slide_root = ET.fromstring(slide_xml)
                    
                    # Look for image references in the slide
                    for elem in slide_root.iter():
                        # Check for blip elements (image references)
                        if elem.tag.endswith('}blip'):
                            for attr_name, attr_value in elem.attrib.items():
                                if 'embed' in attr_name.lower() or 'link' in attr_name.lower():
                                    # This references a media file
                                    media_ref = attr_value
                                    if media_ref not in media_to_slides:
                                        media_to_slides[media_ref] = []
                                    media_to_slides[media_ref].append(slide_num)
                                    print(f"🔗 Found image reference in slide {slide_num}: {media_ref}")
                        
                        # Check for picture elements with position info
                        elif elem.tag.endswith('}pic'):
                            # Look for position and size information in the pic element
                            x = 0
                            y = 0
                            width = 200
                            height = 200
                            
                            # Search for position and size in all child elements
                            for child in elem.iter():
                                if child.tag.endswith('}off'):
                                    # Position information
                                    for attr_name, attr_value in child.attrib.items():
                                        if 'x' in attr_name.lower():
                                            x = int(attr_value) // 9525  # Convert EMU to pixels
                                        elif 'y' in attr_name.lower():
                                            y = int(attr_value) // 9525
                                elif child.tag.endswith('}ext'):
                                    # Size information
                                    for attr_name, attr_value in child.attrib.items():
                                        if 'cx' in attr_name.lower():
                                            width = int(attr_value) // 9525
                                        elif 'cy' in attr_name.lower():
                                            height = int(attr_value) // 9525
                            
                            # Also check for shape properties that might contain position/size
                            for child in elem.iter():
                                if child.tag.endswith('}spPr'):
                                    for grandchild in child.iter():
                                        if grandchild.tag.endswith('}xfrm'):
                                            for attr_name, attr_value in grandchild.attrib.items():
                                                if 'x' in attr_name.lower():
                                                    x = int(attr_value) // 9525
                                                elif 'y' in attr_name.lower():
                                                    y = int(attr_value) // 9525
                                                elif 'cx' in attr_name.lower():
                                                    width = int(attr_value) // 9525
                                                elif 'cy' in attr_name.lower():
                                                    height = int(attr_value) // 9525
                            
                            # Find the corresponding image and update it
                            for img in images:
                                if img.get('slide_index') == -1:  # Unmapped image
                                    img['slide_index'] = slide_num
                                    img['x'] = max(x, 50) if x > 0 else 100  # Default position if not found
                                    img['y'] = max(y, 50) if y > 0 else 100
                                    img['width'] = max(width, 100) if width > 0 else 200  # Default size if not found
                                    img['height'] = max(height, 100) if height > 0 else 150
                                    print(f"✅ Mapped image to slide {slide_num}: pos({img['x']},{img['y']}) size({img['width']}x{img['height']})")
                                    break
                                    
                except Exception as e:
                    print(f"❌ Error parsing slide XML {slide_file}: {e}")
        
        # For any unmapped images, distribute them across slides
        unmapped_images = [img for img in images if img.get('slide_index') == -1]
        if unmapped_images:
            print(f"⚠️ {len(unmapped_images)} images still unmapped, distributing across slides...")
            slides_count = len([f for f in zipfile.ZipFile(pptx_path, 'r').namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')])
            
            for i, img in enumerate(unmapped_images):
                slide_num = i % slides_count
                img['slide_index'] = slide_num
                img['x'] = 100 + (i * 50) % 400  # Spread horizontally
                img['y'] = 100 + (i * 30) % 200  # Spread vertically
                img['width'] = 200
                img['height'] = 150
                print(f"📍 Distributed unmapped image to slide {slide_num}: pos({img['x']},{img['y']})")
        
        print(f"✅ Successfully mapped {len(images)} images to slide positions")
        return images
        
    except Exception as e:
        print(f"❌ Error mapping images to slides: {e}")
        return images

def map_zip_images_to_slides_with_pptx(pptx_path, images, prs):
    """
    Map ZIP-extracted images to their correct slide positions using python-pptx
    This gives us accurate dimensions and positions from the actual slide shapes
    """
    try:
        print(f"🔍 Mapping {len(images)} images using python-pptx for accurate dimensions...")
        
        # Create a list to track which images have been mapped
        mapped_images = []
        image_counter = 0
        
        for slide_num, slide in enumerate(prs.slides):
            print(f"📄 Processing slide {slide_num + 1} with {len(slide.shapes)} shapes...")
            
            # Look for image shapes in this slide
            for shape_idx, shape in enumerate(slide.shapes):
                # Check if this shape is an image - be more comprehensive
                is_image_shape = False
                shape_info = f"Shape {shape_idx}: {type(shape).__name__}"
                
                # Check various ways an image can be stored
                if hasattr(shape, 'image') and shape.image:
                    is_image_shape = True
                    shape_info += " (direct image)"
                elif hasattr(shape, 'fill') and hasattr(shape.fill, 'type') and shape.fill.type in [1, 2, 3]:
                    is_image_shape = True
                    shape_info += f" (fill type {shape.fill.type})"
                elif hasattr(shape, 'image_part') and shape.image_part:
                    is_image_shape = True
                    shape_info += " (image_part)"
                elif type(shape).__name__ == 'Picture':
                    is_image_shape = True
                    shape_info += " (Picture)"
                elif hasattr(shape, 'shape_type') and shape.shape_type == 13:
                    is_image_shape = True
                    shape_info += " (shape_type 13)"
                elif hasattr(shape, 'fill') and hasattr(shape.fill, 'image') and shape.fill.image:
                    is_image_shape = True
                    shape_info += " (fill.image)"
                elif hasattr(shape, 'fill') and hasattr(shape.fill, 'picture') and shape.fill.picture:
                    is_image_shape = True
                    shape_info += " (fill.picture)"
                elif type(shape).__name__ == 'GroupShape':
                    # Check if GroupShape contains image shapes
                    try:
                        if hasattr(shape, 'shapes'):
                            for sub_shape in shape.shapes:
                                if (hasattr(sub_shape, 'image') and sub_shape.image) or \
                                   (hasattr(sub_shape, 'fill') and hasattr(sub_shape.fill, 'type') and sub_shape.fill.type in [1, 2, 3]) or \
                                   (hasattr(sub_shape, 'image_part') and sub_shape.image_part) or \
                                   type(sub_shape).__name__ == 'Picture':
                                    is_image_shape = True
                                    shape_info += " (GroupShape with image)"
                                    break
                    except:
                        pass
                
                # Also check for any shape that might contain image data
                if not is_image_shape:
                    # Check if shape has any image-related attributes
                    for attr_name in dir(shape):
                        if 'image' in attr_name.lower() or 'picture' in attr_name.lower():
                            try:
                                attr_value = getattr(shape, attr_name)
                                if attr_value:
                                    is_image_shape = True
                                    shape_info += f" (has {attr_name})"
                                    break
                            except:
                                pass
                
                if is_image_shape:
                    # Get the actual position and size from the shape
                    x = int(shape.left.inches * 96) if hasattr(shape, 'left') and shape.left else 100
                    y = int(shape.top.inches * 96) if hasattr(shape, 'top') and shape.top else 100
                    width = int(shape.width.inches * 96) if hasattr(shape, 'width') and shape.width else 200
                    height = int(shape.height.inches * 96) if hasattr(shape, 'height') and shape.height else 150
                    
                    print(f"🔍 Found image shape: {shape_info} - pos({x},{y}) size({width}x{height})")
                    
                    # Find an unmapped image to assign to this shape
                    for img in images:
                        if img.get('slide_index') == -1 and img not in mapped_images:
                            img['slide_index'] = slide_num
                            img['x'] = x  # Use exact position
                            img['y'] = y
                            img['width'] = width  # Use exact size
                            img['height'] = height
                            img['shape_index'] = shape_idx
                            
                            mapped_images.append(img)
                            print(f"✅ Mapped image to slide {slide_num}, shape {shape_idx}: pos({img['x']},{img['y']}) size({img['width']}x{img['height']})")
                            break
                else:
                    # Log non-image shapes for debugging
                    if shape_idx < 5:  # Only log first few to avoid spam
                        print(f"📝 Non-image shape: {shape_info}")
        
        # For any remaining unmapped images, distribute them across slides with better positioning
        unmapped_images = [img for img in images if img.get('slide_index') == -1]
        if unmapped_images:
            print(f"⚠️ {len(unmapped_images)} images still unmapped, distributing across slides...")
            slides_count = len(prs.slides)
            
            for i, img in enumerate(unmapped_images):
                slide_num = i % slides_count
                img['slide_index'] = slide_num
                
                # Better distribution - spread images more naturally
                img['x'] = 100 + (i * 80) % 600  # Spread horizontally with more space
                img['y'] = 100 + (i * 60) % 300  # Spread vertically with more space
                img['width'] = 250  # Slightly larger default size
                img['height'] = 200
                print(f"📍 Distributed unmapped image to slide {slide_num}: pos({img['x']},{img['y']}) size({img['width']}x{img['height']})")
        
        print(f"✅ Successfully mapped {len(images)} images to slide positions with accurate dimensions")
        return images
        
    except Exception as e:
        print(f"❌ Error mapping images with python-pptx: {e}")
        import traceback
        traceback.print_exc()
        return images

def extract_all_images_from_pptx(pptx_path):
    """
    Extract all images from PPTX file including:
    - Direct pictures
    - Background fills  
    - Canva blipFill images
    - Linked images (URLs)
    
    Returns JSON objects with slide_index, shape_index, src, x, y, width, height
    """
    images = []
    image_counter = 0
    
    try:
        prs = Presentation(pptx_path)
        
        # Also check slide relationships for linked images
        slide_relationships = {}
        try:
            for slide_idx, slide in enumerate(prs.slides):
                if hasattr(slide, 'part') and hasattr(slide.part, 'rels'):
                    slide_relationships[slide_idx] = slide.part.rels
        except Exception as e:
            print(f"⚠️ Could not extract slide relationships: {e}")
        
        for slide_idx, slide in enumerate(prs.slides):
            print(f"🔍 Processing slide {slide_idx + 1} with {len(slide.shapes)} shapes")
            
            for shape_idx, shape in enumerate(slide.shapes):
                print(f"  🔍 Checking shape {shape_idx}: {type(shape).__name__}")
                
                # Debug: Print detailed shape information for Canva debugging
                if hasattr(shape, 'fill'):
                    print(f"    📋 Fill type: {getattr(shape.fill, 'type', 'unknown')}")
                    if hasattr(shape.fill, 'image'):
                        print(f"    🖼️ Has fill.image: {shape.fill.image is not None}")
                    if hasattr(shape.fill, 'picture'):
                        print(f"    🖼️ Has fill.picture: {shape.fill.picture is not None}")
                    if hasattr(shape.fill, 'blipFill'):
                        print(f"    🖼️ Has fill.blipFill: {shape.fill.blipFill is not None}")
                if hasattr(shape, 'image'):
                    print(f"    🖼️ Has direct image: {shape.image is not None}")
                if hasattr(shape, 'image_part'):
                    print(f"    🖼️ Has image_part: {shape.image_part is not None}")
                if hasattr(shape, 'shape_type'):
                    print(f"    📐 Shape type: {shape.shape_type}")
                
                # Comprehensive image detection - try multiple methods
                image_found = False
                
                # Method 1: Direct image attribute
                if hasattr(shape, 'image') and shape.image:
                    try:
                        image_data = shape.image.blob
                        if image_data and len(image_data) > 0:
                            image_base64 = base64.b64encode(image_data).decode('utf-8')
                            
                            # Determine MIME type
                            mime_type = 'image/jpeg'
                            if image_data.startswith(b'\x89PNG'):
                                mime_type = 'image/png'
                            elif image_data.startswith(b'GIF'):
                                mime_type = 'image/gif'
                            elif image_data.startswith(b'RIFF') and b'WEBP' in image_data[:12]:
                                mime_type = 'image/webp'
                            
                            data_url = f"data:{mime_type};base64,{image_base64}"
                            image_id = f"image_{slide_num}_{shape_idx}_{image_counter}"
                            images[image_id] = data_url
                            image_counter += 1
                            
                            print(f"✅ Extracted image: {image_id} ({len(image_data)} bytes, {mime_type})")
                            image_found = True
                        else:
                            print(f"⚠️ Image data is empty for shape {shape_idx}")
                    except Exception as e:
                        print(f"❌ Failed to extract image from shape {shape_idx}: {e}")
                
                # Method 2: Try all possible fill types and attributes
                if not image_found:
                    print(f"🔍 Trying alternative image detection methods for shape {shape_idx}")
                    
                    # Try different fill types
                    if hasattr(shape, 'fill'):
                        print(f"    📋 Fill type: {getattr(shape.fill, 'type', 'unknown')}")
                        
                        # Try fill.image
                        if hasattr(shape.fill, 'image') and shape.fill.image:
                            try:
                                image_data = shape.fill.image.blob
                                if image_data and len(image_data) > 0:
                                    image_base64 = base64.b64encode(image_data).decode('utf-8')
                                    
                                    mime_type = 'image/jpeg'
                                    if image_data.startswith(b'\x89PNG'):
                                        mime_type = 'image/png'
                                    elif image_data.startswith(b'GIF'):
                                        mime_type = 'image/gif'
                                    
                                    data_url = f"data:{mime_type};base64,{image_base64}"
                                    image_id = f"image_{slide_num}_{shape_idx}_{image_counter}"
                                    images[image_id] = data_url
                                    image_counter += 1
                                    
                                    print(f"✅ Extracted image from fill.image: {image_id} ({len(image_data)} bytes, {mime_type})")
                                    image_found = True
                            except Exception as e:
                                print(f"❌ Failed to extract from fill.image: {e}")
                        
                        # Try fill.picture
                        if not image_found and hasattr(shape.fill, 'picture') and shape.fill.picture:
                            try:
                                image_data = shape.fill.picture.blob
                                if image_data and len(image_data) > 0:
                                    image_base64 = base64.b64encode(image_data).decode('utf-8')
                                    
                                    mime_type = 'image/jpeg'
                                    if image_data.startswith(b'\x89PNG'):
                                        mime_type = 'image/png'
                                    elif image_data.startswith(b'GIF'):
                                        mime_type = 'image/gif'
                                    
                                    data_url = f"data:{mime_type};base64,{image_base64}"
                                    image_id = f"image_{slide_num}_{shape_idx}_{image_counter}"
                                    images[image_id] = data_url
                                    image_counter += 1
                                    
                                    print(f"✅ Extracted image from fill.picture: {image_id} ({len(image_data)} bytes, {mime_type})")
                                    image_found = True
                            except Exception as e:
                                print(f"❌ Failed to extract from fill.picture: {e}")
                        
                        # Try fill.blipFill (PowerPoint format)
                        if not image_found and hasattr(shape.fill, 'blipFill') and shape.fill.blipFill:
                            try:
                                if hasattr(shape.fill.blipFill, 'blip') and shape.fill.blipFill.blip:
                                    if hasattr(shape.fill.blipFill.blip, 'blob'):
                                        image_data = shape.fill.blipFill.blip.blob
                                        if image_data and len(image_data) > 0:
                                            image_base64 = base64.b64encode(image_data).decode('utf-8')
                                            
                                            mime_type = 'image/jpeg'
                                            if image_data.startswith(b'\x89PNG'):
                                                mime_type = 'image/png'
                                            elif image_data.startswith(b'GIF'):
                                                mime_type = 'image/gif'
                                            
                                            data_url = f"data:{mime_type};base64,{image_base64}"
                                            image_id = f"image_{slide_num}_{shape_idx}_{image_counter}"
                                            images[image_id] = data_url
                                            image_counter += 1
                                            
                                            print(f"✅ Extracted image from fill.blipFill: {image_id} ({len(image_data)} bytes, {mime_type})")
                                            image_found = True
                            except Exception as e:
                                print(f"❌ Failed to extract from fill.blipFill: {e}")
                    
                    # Method 3: Try image_part
                    if not image_found and hasattr(shape, 'image_part') and shape.image_part:
                        try:
                            image_data = shape.image_part.blob
                            if image_data and len(image_data) > 0:
                                image_base64 = base64.b64encode(image_data).decode('utf-8')
                                
                                mime_type = 'image/jpeg'
                                if image_data.startswith(b'\x89PNG'):
                                    mime_type = 'image/png'
                                elif image_data.startswith(b'GIF'):
                                    mime_type = 'image/gif'
                                
                                data_url = f"data:{mime_type};base64,{image_base64}"
                                image_id = f"image_{slide_num}_{shape_idx}_{image_counter}"
                                images[image_id] = data_url
                                image_counter += 1
                                
                                print(f"✅ Extracted image from image_part: {image_id} ({len(image_data)} bytes, {mime_type})")
                                image_found = True
                        except Exception as e:
                            print(f"❌ Failed to extract from image_part: {e}")
                    
                    if not image_found:
                        print(f"❌ No image found in shape {shape_idx} using any method")
                        
                        # Method 4: Deep inspection of shape attributes
                        print(f"🔍 Deep inspection of shape {shape_idx} attributes:")
                        for attr_name in dir(shape):
                            if not attr_name.startswith('_'):
                                try:
                                    attr_value = getattr(shape, attr_name)
                                    if hasattr(attr_value, 'blob') or 'image' in attr_name.lower():
                                        print(f"    🔍 Found potential image attribute: {attr_name} = {type(attr_value)}")
                                        if hasattr(attr_value, 'blob'):
                                            try:
                                                blob_data = attr_value.blob
                                                if blob_data and len(blob_data) > 0:
                                                    print(f"    ✅ Found blob data in {attr_name}: {len(blob_data)} bytes")
                                                    image_base64 = base64.b64encode(blob_data).decode('utf-8')
                                                    
                                                    mime_type = 'image/jpeg'
                                                    if blob_data.startswith(b'\x89PNG'):
                                                        mime_type = 'image/png'
                                                    elif blob_data.startswith(b'GIF'):
                                                        mime_type = 'image/gif'
                                                    
                                                    data_url = f"data:{mime_type};base64,{image_base64}"
                                                    image_id = f"image_{slide_num}_{shape_idx}_{image_counter}"
                                                    images[image_id] = data_url
                                                    image_counter += 1
                                                    
                                                    print(f"✅ Extracted image from {attr_name}: {image_id} ({len(blob_data)} bytes, {mime_type})")
                                                    image_found = True
                                                    break
                                            except Exception as e:
                                                print(f"    ❌ Error accessing blob in {attr_name}: {e}")
                                except Exception as e:
                                    pass  # Skip attributes that can't be accessed
                
                # Check for PICTURE fill (type 3 = MSO_FILL_PICTURE)
                elif hasattr(shape, 'fill') and hasattr(shape.fill, 'type') and shape.fill.type == 3:
                    try:
                        if hasattr(shape.fill, 'image') and shape.fill.image:
                            image_data = shape.fill.image.blob
                            if image_data and len(image_data) > 0:
                                image_base64 = base64.b64encode(image_data).decode('utf-8')
                                
                                # Determine MIME type
                                mime_type = 'image/jpeg'
                                if image_data.startswith(b'\x89PNG'):
                                    mime_type = 'image/png'
                                elif image_data.startswith(b'GIF'):
                                    mime_type = 'image/gif'
                                elif image_data.startswith(b'RIFF') and b'WEBP' in image_data[:12]:
                                    mime_type = 'image/webp'
                                
                                data_url = f"data:{mime_type};base64,{image_base64}"
                                image_id = f"image_{slide_num}_{shape_idx}_{image_counter}"
                                images[image_id] = data_url
                                image_counter += 1
                                
                                print(f"✅ Extracted fill image: {image_id} ({len(image_data)} bytes, {mime_type})")
                            else:
                                print(f"⚠️ Fill image data is empty for shape {shape_idx}")
                    except Exception as e:
                        print(f"❌ Failed to extract fill image from shape {shape_idx}: {e}")
                
                # Check for other image types
                elif hasattr(shape, 'image_part') and shape.image_part:
                    try:
                        image_data = shape.image_part.blob
                        if image_data and len(image_data) > 0:
                            image_base64 = base64.b64encode(image_data).decode('utf-8')
                            
                            # Determine MIME type
                            mime_type = 'image/jpeg'
                            if image_data.startswith(b'\x89PNG'):
                                mime_type = 'image/png'
                            elif image_data.startswith(b'GIF'):
                                mime_type = 'image/gif'
                            
                            data_url = f"data:{mime_type};base64,{image_base64}"
                            image_id = f"image_{slide_num}_{shape_idx}_{image_counter}"
                            images[image_id] = data_url
                            image_counter += 1
                            
                            print(f"✅ Extracted image_part: {image_id} ({len(image_data)} bytes)")
                    except Exception as e:
                        print(f"❌ Failed to extract image_part from shape {shape_idx}: {e}")
    
    except Exception as e:
        print(f"❌ Error processing PPTX: {e}")
    
    print(f"📁 Total images extracted: {len(images)}")
    return images

def extract_images_from_zip_structure(pptx_path):
    """
    Extract images directly from PPTX ZIP structure (like Pages does)
    This method reads the raw ZIP file and extracts images from ppt/media/
    """
    images = []
    image_counter = 0
    
    try:
        print(f"🔍 Extracting images from ZIP structure of {pptx_path}")
        
        with zipfile.ZipFile(pptx_path, 'r') as zip_file:
            # List all files in the ZIP
            file_list = zip_file.namelist()
            print(f"📁 Found {len(file_list)} files in ZIP structure")
            
            # Find all media files
            media_files = [f for f in file_list if f.startswith('ppt/media/')]
            print(f"🖼️ Found {len(media_files)} media files: {media_files}")
            
            # Extract each media file
            for media_file in media_files:
                try:
                    # Read the image data
                    image_data = zip_file.read(media_file)
                    
                    if image_data and len(image_data) > 0:
                        # Determine file extension and MIME type
                        file_ext = media_file.split('.')[-1].lower()
                        mime_type = 'image/jpeg'
                        
                        if file_ext == 'png':
                            mime_type = 'image/png'
                        elif file_ext == 'gif':
                            mime_type = 'image/gif'
                        elif file_ext == 'webp':
                            mime_type = 'image/webp'
                        elif file_ext == 'bmp':
                            mime_type = 'image/bmp'
                        elif file_ext == 'tiff':
                            mime_type = 'image/tiff'
                        
                        # Compress image to reduce size for Firebase
                        compressed_data, size_kb = compress_image(image_data, max_size_kb=80, quality=75)
                        
                        # Convert to base64 (always use JPEG after compression)
                        image_base64 = base64.b64encode(compressed_data).decode('utf-8')
                        image_src = f"data:image/jpeg;base64,{image_base64}"
                        
                        # Create image object with placeholder position (will be updated later)
                        image_obj = {
                            "slide_index": -1,  # Will be mapped to actual slide
                            "shape_index": image_counter,
                            "src": image_src,
                            "x": 0,
                            "y": 0,
                            "width": 200,
                            "height": 200,
                            "filename": media_file,
                            "size_bytes": len(compressed_data),
                            "original_size_kb": len(image_data) / 1024,
                            "compressed_size_kb": size_kb,
                            "media_file": media_file  # Store original filename for mapping
                        }
                        
                        images.append(image_obj)
                        image_counter += 1
                        
                        print(f"✅ Extracted image from ZIP: {media_file} ({len(image_data)} bytes, {mime_type})")
                        
                except Exception as e:
                    print(f"❌ Error extracting {media_file}: {e}")
            
            # Now try to map images to slides by reading slide XML files
            slide_files = [f for f in file_list if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
            print(f"📄 Found {len(slide_files)} slide XML files")
            
            # Parse slide relationships to map images to slides
            for slide_file in slide_files:
                try:
                    slide_xml = zip_file.read(slide_file)
                    slide_root = ET.fromstring(slide_xml)
                    
                    # Extract slide number from filename (handle both slide1.xml and slides/slide1.xml)
                    if 'slides/' in slide_file:
                        slide_num = int(slide_file.split('slides/slide')[1].split('.')[0]) - 1
                    else:
                        slide_num = int(slide_file.split('slide')[1].split('.')[0]) - 1
                    
                    # Look for image references in the slide XML
                    for elem in slide_root.iter():
                        if elem.tag.endswith('}blip') or elem.tag.endswith('}pic'):
                            # This might reference an image
                            for attr_name, attr_value in elem.attrib.items():
                                if 'embed' in attr_name.lower() or 'link' in attr_name.lower():
                                    print(f"🔗 Found image reference in slide {slide_num}: {attr_name}={attr_value}")
                                    
                except Exception as e:
                    print(f"❌ Error parsing slide XML {slide_file}: {e}")
        
        print(f"📊 Total images extracted from ZIP: {len(images)}")
        return images
        
    except Exception as e:
        print(f"❌ Error extracting images from ZIP structure: {e}")
        return []

def extract_all_images_comprehensive(pptx_path):
    """
    Comprehensive image extraction function that works with Canva and all PPTX formats.
    Returns JSON objects with slide_index, shape_index, src, x, y, width, height
    """
    images = []
    image_counter = 0
    
    try:
        prs = Presentation(pptx_path)
        
        # Check slide relationships for linked images
        slide_relationships = {}
        try:
            for slide_idx, slide in enumerate(prs.slides):
                if hasattr(slide, 'part') and hasattr(slide.part, 'rels'):
                    slide_relationships[slide_idx] = slide.part.rels
        except Exception as e:
            print(f"⚠️ Could not extract slide relationships: {e}")
        
        for slide_idx, slide in enumerate(prs.slides):
            print(f"🔍 Processing slide {slide_idx + 1} with {len(slide.shapes)} shapes")
            
            for shape_idx, shape in enumerate(slide.shapes):
                print(f"  🔍 Checking shape {shape_idx}: {type(shape).__name__}")
                
                image_found = False
                image_src = ""
                
                # Method 1: Direct image attribute
                if hasattr(shape, 'image') and shape.image:
                    try:
                        image_data = shape.image.blob
                        if image_data and len(image_data) > 0:
                            image_base64 = base64.b64encode(image_data).decode('utf-8')
                            
                            mime_type = 'image/jpeg'
                            if image_data.startswith(b'\x89PNG'):
                                mime_type = 'image/png'
                            elif image_data.startswith(b'GIF'):
                                mime_type = 'image/gif'
                            elif image_data.startswith(b'RIFF') and b'WEBP' in image_data[:12]:
                                mime_type = 'image/webp'
                            
                            image_src = f"data:{mime_type};base64,{image_base64}"
                            print(f"✅ Extracted direct image from shape {shape_idx} ({len(image_data)} bytes, {mime_type})")
                            image_found = True
                    except Exception as e:
                        print(f"❌ Failed to extract direct image: {e}")
                
                # Method 2: Fill-based images (Canva often uses these)
                if not image_found and hasattr(shape, 'fill'):
                    # Try fill.image
                    if hasattr(shape.fill, 'image') and shape.fill.image:
                        try:
                            image_data = shape.fill.image.blob
                            if image_data and len(image_data) > 0:
                                image_base64 = base64.b64encode(image_data).decode('utf-8')
                                
                                mime_type = 'image/jpeg'
                                if image_data.startswith(b'\x89PNG'):
                                    mime_type = 'image/png'
                                elif image_data.startswith(b'GIF'):
                                    mime_type = 'image/gif'
                                
                                image_src = f"data:{mime_type};base64,{image_base64}"
                                print(f"✅ Extracted fill.image from shape {shape_idx} ({len(image_data)} bytes, {mime_type})")
                                image_found = True
                        except Exception as e:
                            print(f"❌ Failed to extract from fill.image: {e}")
                    
                    # Try fill.picture
                    if not image_found and hasattr(shape.fill, 'picture') and shape.fill.picture:
                        try:
                            image_data = shape.fill.picture.blob
                            if image_data and len(image_data) > 0:
                                image_base64 = base64.b64encode(image_data).decode('utf-8')
                                
                                mime_type = 'image/jpeg'
                                if image_data.startswith(b'\x89PNG'):
                                    mime_type = 'image/png'
                                elif image_data.startswith(b'GIF'):
                                    mime_type = 'image/gif'
                                
                                image_src = f"data:{mime_type};base64,{image_base64}"
                                print(f"✅ Extracted fill.picture from shape {shape_idx} ({len(image_data)} bytes, {mime_type})")
                                image_found = True
                        except Exception as e:
                            print(f"❌ Failed to extract from fill.picture: {e}")
                    
                    # Try fill.blipFill (PowerPoint and Canva format)
                    if not image_found and hasattr(shape.fill, 'blipFill') and shape.fill.blipFill:
                        try:
                            if hasattr(shape.fill.blipFill, 'blip') and shape.fill.blipFill.blip:
                                if hasattr(shape.fill.blipFill.blip, 'blob'):
                                    image_data = shape.fill.blipFill.blip.blob
                                    if image_data and len(image_data) > 0:
                                        image_base64 = base64.b64encode(image_data).decode('utf-8')
                                        
                                        mime_type = 'image/jpeg'
                                        if image_data.startswith(b'\x89PNG'):
                                            mime_type = 'image/png'
                                        elif image_data.startswith(b'GIF'):
                                            mime_type = 'image/gif'
                                        
                                        image_src = f"data:{mime_type};base64,{image_base64}"
                                        print(f"✅ Extracted fill.blipFill from shape {shape_idx} ({len(image_data)} bytes, {mime_type})")
                                        image_found = True
                        except Exception as e:
                            print(f"❌ Failed to extract from fill.blipFill: {e}")
                
                # Method 3: image_part
                if not image_found and hasattr(shape, 'image_part') and shape.image_part:
                    try:
                        image_data = shape.image_part.blob
                        if image_data and len(image_data) > 0:
                            image_base64 = base64.b64encode(image_data).decode('utf-8')
                            
                            mime_type = 'image/jpeg'
                            if image_data.startswith(b'\x89PNG'):
                                mime_type = 'image/png'
                            elif image_data.startswith(b'GIF'):
                                mime_type = 'image/gif'
                            
                            image_src = f"data:{mime_type};base64,{image_base64}"
                            print(f"✅ Extracted image_part from shape {shape_idx} ({len(image_data)} bytes, {mime_type})")
                            image_found = True
                    except Exception as e:
                        print(f"❌ Failed to extract from image_part: {e}")
                
                # Method 4: Deep attribute inspection
                if not image_found:
                    for attr_name in dir(shape):
                        if not attr_name.startswith('_'):
                            try:
                                attr_value = getattr(shape, attr_name)
                                if hasattr(attr_value, 'blob') or 'image' in attr_name.lower():
                                    if hasattr(attr_value, 'blob'):
                                        try:
                                            blob_data = attr_value.blob
                                            if blob_data and len(blob_data) > 0:
                                                image_base64 = base64.b64encode(blob_data).decode('utf-8')
                                                
                                                mime_type = 'image/jpeg'
                                                if blob_data.startswith(b'\x89PNG'):
                                                    mime_type = 'image/png'
                                                elif blob_data.startswith(b'GIF'):
                                                    mime_type = 'image/gif'
                                                
                                                image_src = f"data:{mime_type};base64,{image_base64}"
                                                print(f"✅ Extracted image from {attr_name}: shape {shape_idx} ({len(blob_data)} bytes, {mime_type})")
                                                image_found = True
                                                break
                                        except Exception as e:
                                            pass
                            except Exception as e:
                                pass
                
                # Method 5: Check for linked images (URLs)
                if not image_found and slide_idx in slide_relationships:
                    try:
                        for rel in slide_relationships[slide_idx].values():
                            if hasattr(rel, 'target_ref') and rel.target_ref:
                                if any(img_ext in rel.target_ref.lower() for img_ext in ['.jpg', '.jpeg', '.png', '.gif', '.webp']):
                                    image_src = rel.target_ref
                                    print(f"✅ Found linked image: {image_src}")
                                    image_found = True
                                    break
                    except Exception as e:
                        print(f"❌ Error checking linked images: {e}")
                
                # Create JSON object if image found
                if image_found and image_src:
                    image_obj = {
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "src": image_src,
                        "x": int(shape.left.inches * 96) if hasattr(shape, 'left') else 0,
                        "y": int(shape.top.inches * 96) if hasattr(shape, 'top') else 0,
                        "width": int(shape.width.inches * 96) if hasattr(shape, 'width') else 200,
                        "height": int(shape.height.inches * 96) if hasattr(shape, 'height') else 200
                    }
                    images.append(image_obj)
                    image_counter += 1
                    print(f"📸 Added image {image_counter}: slide {slide_idx}, shape {shape_idx}, size {image_obj['width']}x{image_obj['height']}")
                else:
                    print(f"❌ No image found in shape {shape_idx}")
        
        print(f"📊 Total images extracted: {len(images)}")
        return images
        
    except Exception as e:
        print(f"❌ Error extracting images from PPTX: {e}")
        return []

def parse_pptx_to_json(pptx_path):
    """Parse PPTX file and return structured JSON"""
    try:
        prs = Presentation(pptx_path)
        # Try ZIP extraction first (like Pages does)
        images = extract_images_from_zip_structure(pptx_path)
        
        # If no images found via ZIP, try comprehensive extraction
        if not images:
            print("🔄 No images found via ZIP extraction, trying comprehensive method...")
            images = extract_all_images_comprehensive(pptx_path)
        
        # Map ZIP images to their correct slide positions using python-pptx
        if images:
            print(f"🗺️ Mapping {len(images)} ZIP images to slide positions...")
            images = map_zip_images_to_slides_with_pptx(pptx_path, images, prs)
        
        slides = []
        image_index = 0
        
        for slide_num, slide in enumerate(prs.slides):
            # Extract background color from slide
            background_color = "#ffffff"  # Default white
            try:
                if hasattr(slide, 'background') and slide.background:
                    if hasattr(slide.background, 'fill') and slide.background.fill:
                        if hasattr(slide.background.fill, 'solid_color') and slide.background.fill.solid_color:
                            color = slide.background.fill.solid_color
                            if hasattr(color, 'rgb') and color.rgb:
                                rgb = color.rgb
                                if hasattr(rgb, 'red') and hasattr(rgb, 'green') and hasattr(rgb, 'blue'):
                                    background_color = f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
                                    print(f"🎨 Extracted slide background color: {background_color}")
                        elif hasattr(slide.background.fill, 'fore_color') and slide.background.fill.fore_color:
                            color = slide.background.fill.fore_color
                            if hasattr(color, 'rgb') and color.rgb:
                                rgb = color.rgb
                                if hasattr(rgb, 'red') and hasattr(rgb, 'green') and hasattr(rgb, 'blue'):
                                    background_color = f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
                                    print(f"🎨 Extracted slide background color (fore): {background_color}")
            except Exception as e:
                print(f"⚠️ Could not extract background color: {e}")
            
            slide_data = {
                "id": f"slide-{slide_num + 1}",
                "title": f"Slide {slide_num + 1}",
                "content": "",
                "elements": [],
                "background": background_color
            }
            
            # Process all shapes in the slide
            for shape_idx, shape in enumerate(slide.shapes):
                # Handle text shapes
                if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text.strip():
                    text_content = shape.text_frame.text.strip()
                    
                    # Extract font information from the first paragraph
                    font_size = 24
                    font_family = "Inter"
                    font_weight = "600"
                    text_color = "#000000"
                    text_align = "left"
                    
                    try:
                        if shape.text_frame.paragraphs:
                            first_para = shape.text_frame.paragraphs[0]
                            if first_para.runs:
                                first_run = first_para.runs[0]
                                
                                # Extract font size
                                if hasattr(first_run.font, 'size') and first_run.font.size:
                                    font_size = int(first_run.font.size.pt)
                                
                                # Extract font family
                                if hasattr(first_run.font, 'name') and first_run.font.name:
                                    font_family = first_run.font.name
                                
                                # Extract font weight (bold)
                                if hasattr(first_run.font, 'bold') and first_run.font.bold:
                                    font_weight = "bold"
                                
                                # Extract italic style
                                if hasattr(first_run.font, 'italic') and first_run.font.italic:
                                    if font_family and 'italic' not in font_family.lower():
                                        font_family = f"{font_family} Italic"
                                
                                # Extract text color
                                if hasattr(first_run.font, 'color') and first_run.font.color:
                                    try:
                                        print(f"🎨 Extracting color from shape {shape_idx}")
                                        print(f"🎨 Color object: {first_run.font.color}")
                                        print(f"🎨 Has rgb: {hasattr(first_run.font.color, 'rgb')}")
                                        
                                        if hasattr(first_run.font.color, 'rgb'):
                                            rgb = first_run.font.color.rgb
                                            print(f"🎨 RGB object: {rgb}")
                                            if rgb:
                                                # Handle different RGB color formats
                                                if hasattr(rgb, 'red') and hasattr(rgb, 'green') and hasattr(rgb, 'blue'):
                                                    # Standard RGB format
                                                    text_color = f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
                                                    print(f"🎨 Extracted color (standard): {text_color}")
                                                elif hasattr(rgb, 'r') and hasattr(rgb, 'g') and hasattr(rgb, 'b'):
                                                    # Alternative RGB format
                                                    text_color = f"#{rgb.r:02x}{rgb.g:02x}{rgb.b:02x}"
                                                    print(f"🎨 Extracted color (alt): {text_color}")
                                                else:
                                                    # Try to get RGB values directly from the string representation
                                                    rgb_values = str(rgb)
                                                    print(f"🎨 RGB string: {rgb_values}")
                                                    
                                                    # Handle direct hex values like "422717"
                                                    if len(rgb_values) == 6 and all(c in '0123456789abcdefABCDEF' for c in rgb_values):
                                                        text_color = f"#{rgb_values}"
                                                        print(f"🎨 Extracted color (direct hex): {text_color}")
                                                    elif 'RGB' in rgb_values:
                                                        # Extract RGB values from string representation
                                                        import re
                                                        match = re.search(r'RGB\((\d+),\s*(\d+),\s*(\d+)\)', rgb_values)
                                                        if match:
                                                            r, g, b = map(int, match.groups())
                                                            text_color = f"#{r:02x}{g:02x}{b:02x}"
                                                            print(f"🎨 Extracted color (regex): {text_color}")
                                        elif hasattr(first_run.font.color, 'theme_color'):
                                            # Handle theme colors - use a default color for now
                                            text_color = "#000000"
                                            print(f"🎨 Using default color for theme: {text_color}")
                                    except Exception as color_error:
                                        print(f"⚠️ Color extraction error: {color_error}")
                                        text_color = "#000000"
                                else:
                                    print(f"🎨 No color found for shape {shape_idx}")
                                
                                # Extract text alignment
                                if hasattr(first_para, 'alignment'):
                                    alignment_map = {
                                        1: "left",    # PP_ALIGN_LEFT
                                        2: "center",  # PP_ALIGN_CENTER
                                        3: "right",   # PP_ALIGN_RIGHT
                                        4: "justify"  # PP_ALIGN_JUSTIFY
                                    }
                                    text_align = alignment_map.get(first_para.alignment, "left")
                                    
                    except Exception as e:
                        print(f"⚠️ Could not extract font info: {e}")
                    
                    element = {
                        "id": f"text-{slide_num}-{shape_idx}",
                        "type": "text",
                        "x": int(shape.left.inches * 96),
                        "y": int(shape.top.inches * 96),
                        "width": max(int(shape.width.inches * 96), 200),
                        "height": max(int(shape.height.inches * 96), 60),
                        "content": text_content,
                        "fontSize": font_size,
                        "fontFamily": font_family,
                        "fontWeight": font_weight,
                        "color": text_color,
                        "textAlign": text_align,
                        "rotation": 0,
                        "zIndex": 1,
                        "selected": False
                    }
                    slide_data["elements"].append(element)
                    print(f"✅ Created text element: {text_content[:50]}... (font: {font_family}, size: {font_size}, color: {text_color})")
                
                # Handle image shapes - use ZIP extraction results first
                elif (hasattr(shape, 'image') and shape.image) or \
                     (hasattr(shape, 'fill') and hasattr(shape.fill, 'type') and shape.fill.type in [1, 2, 3]) or \
                     (hasattr(shape, 'image_part') and shape.image_part) or \
                     type(shape).__name__ == 'Picture' or \
                     (type(shape).__name__ == 'GroupShape' and hasattr(shape, 'shapes') and any(
                         (hasattr(sub_shape, 'image') and sub_shape.image) or 
                         (hasattr(sub_shape, 'fill') and hasattr(sub_shape.fill, 'type') and sub_shape.fill.type in [1, 2, 3]) or
                         (hasattr(sub_shape, 'image_part') and sub_shape.image_part) or
                         type(sub_shape).__name__ == 'Picture'
                         for sub_shape in shape.shapes
                     )):
                    # Initialize image_src variable
                    image_src = ""
                    
                    # Check if this shape corresponds to an image from our ZIP extraction
                    matching_image = None
                    for img in images:
                        if img['slide_index'] == slide_num and img['shape_index'] == shape_idx:
                            matching_image = img
                            break
                    
                    if matching_image:
                        print(f"🖼️ Found matching ZIP image for shape {shape_idx}: {type(shape).__name__}")
                        
                        element = {
                            "id": f"image-{slide_num}-{shape_idx}",
                            "type": "image",
                            "x": matching_image['x'],
                            "y": matching_image['y'],
                            "width": matching_image['width'],
                            "height": matching_image['height'],
                            "src": matching_image['src'],
                            "alt": f"Image from slide {slide_num + 1}",
                            "rotation": 0,
                            "zIndex": 1,
                            "selected": False
                        }
                        slide_data["elements"].append(element)
                        print(f"✅ Created ZIP image element: {element['id']} (size: {element['width']}x{element['height']})")
                        
                        # Debug: Print image src preview
                        print(f"🔍 ZIP Image src preview: {matching_image['src'][:100]}...")
                        print(f"🔍 Full ZIP image element: {element}")
                    else:
                        # Check if this might be an image shape that wasn't caught by ZIP extraction
                        is_potential_image = False
                        if (hasattr(shape, 'image') and shape.image) or \
                           (hasattr(shape, 'fill') and hasattr(shape.fill, 'type') and shape.fill.type in [1, 2, 3]) or \
                           (hasattr(shape, 'image_part') and shape.image_part) or \
                           (type(shape).__name__ == 'Picture') or \
                           (hasattr(shape, 'shape_type') and shape.shape_type == 13):
                            is_potential_image = True
                        
                        if is_potential_image:
                            print(f"⚠️ Potential image shape {shape_idx} not found in ZIP extraction: {type(shape).__name__}")
                            
                            # Try to create a placeholder image element with the shape's dimensions
                            try:
                                x = int(shape.left.inches * 96) if hasattr(shape, 'left') and shape.left else 100
                                y = int(shape.top.inches * 96) if hasattr(shape, 'top') and shape.top else 100
                                width = int(shape.width.inches * 96) if hasattr(shape, 'width') and shape.width else 200
                                height = int(shape.height.inches * 96) if hasattr(shape, 'height') and shape.height else 150
                                
                                # Create a placeholder image element
                                element = {
                                    "id": f"image-{slide_num}-{shape_idx}",
                                    "type": "image",
                                    "x": x,
                                    "y": y,
                                    "width": width,
                                    "height": height,
                                    "src": "data:image/svg+xml;base64,PHN2ZyB3aWR0aD0iMjAwIiBoZWlnaHQ9IjE1MCIgeG1sbnM9Imh0dHA6Ly93d3cudzMub3JnLzIwMDAvc3ZnIj48cmVjdCB3aWR0aD0iMTAwJSIgaGVpZ2h0PSIxMDAlIiBmaWxsPSIjZGRkIi8+PHRleHQgeD0iNTAlIiB5PSI1MCUiIGZvbnQtZmFtaWx5PSJBcmlhbCIgZm9udC1zaXplPSIxNCIgZmlsbD0iIzk5OSIgdGV4dC1hbmNob3I9Im1pZGRsZSIgZHk9Ii4zZW0iPkltYWdlPC90ZXh0Pjwvc3ZnPg==",
                                    "alt": f"Placeholder image from slide {slide_num + 1}",
                                    "rotation": 0,
                                    "zIndex": 1,
                                    "selected": False
                                }
                                slide_data["elements"].append(element)
                                print(f"✅ Created placeholder image element: {element['id']} (size: {element['width']}x{element['height']})")
                            except Exception as e:
                                print(f"❌ Error creating placeholder image: {e}")
                    try:
                        # Method 1: Direct image attribute
                        if hasattr(shape, 'image') and shape.image:
                            image_data = shape.image.blob
                            if image_data and len(image_data) > 0:
                                image_base64 = base64.b64encode(image_data).decode('utf-8')
                                
                                # Determine MIME type
                                mime_type = 'image/jpeg'
                                if image_data.startswith(b'\x89PNG'):
                                    mime_type = 'image/png'
                                elif image_data.startswith(b'GIF'):
                                    mime_type = 'image/gif'
                                elif image_data.startswith(b'RIFF') and b'WEBP' in image_data[:12]:
                                    mime_type = 'image/webp'
                                
                                image_src = f"data:{mime_type};base64,{image_base64}"
                                print(f"✅ Extracted image directly from shape {shape_idx} ({len(image_data)} bytes, {mime_type})")
                            else:
                                print(f"⚠️ Image data is empty for shape {shape_idx}")
                        
                        # Method 2: Try different fill types (Canva often uses these)
                        elif hasattr(shape, 'fill') and hasattr(shape.fill, 'type'):
                            print(f"🔍 Checking fill type {shape.fill.type} for shape {shape_idx}")
                            
                            # Try fill.image
                            if hasattr(shape.fill, 'image') and shape.fill.image:
                                image_data = shape.fill.image.blob
                                if image_data and len(image_data) > 0:
                                    image_base64 = base64.b64encode(image_data).decode('utf-8')
                                    
                                    # Determine MIME type
                                    mime_type = 'image/jpeg'
                                    if image_data.startswith(b'\x89PNG'):
                                        mime_type = 'image/png'
                                    elif image_data.startswith(b'GIF'):
                                        mime_type = 'image/gif'
                                    
                                    image_src = f"data:{mime_type};base64,{image_base64}"
                                    print(f"✅ Extracted fill.image from shape {shape_idx} ({len(image_data)} bytes, {mime_type})")
                            
                            # Try fill.picture
                            elif hasattr(shape.fill, 'picture') and shape.fill.picture:
                                image_data = shape.fill.picture.blob
                                if image_data and len(image_data) > 0:
                                    image_base64 = base64.b64encode(image_data).decode('utf-8')
                                    
                                    # Determine MIME type
                                    mime_type = 'image/jpeg'
                                    if image_data.startswith(b'\x89PNG'):
                                        mime_type = 'image/png'
                                    elif image_data.startswith(b'GIF'):
                                        mime_type = 'image/gif'
                                    
                                    image_src = f"data:{mime_type};base64,{image_base64}"
                                    print(f"✅ Extracted fill.picture from shape {shape_idx} ({len(image_data)} bytes, {mime_type})")
                            
                            # Try fill.blipFill (common in PowerPoint)
                            elif hasattr(shape.fill, 'blipFill') and shape.fill.blipFill:
                                if hasattr(shape.fill.blipFill, 'blip') and shape.fill.blipFill.blip:
                                    if hasattr(shape.fill.blipFill.blip, 'blob'):
                                        image_data = shape.fill.blipFill.blip.blob
                                        if image_data and len(image_data) > 0:
                                            image_base64 = base64.b64encode(image_data).decode('utf-8')
                                            
                                            # Determine MIME type
                                            mime_type = 'image/jpeg'
                                            if image_data.startswith(b'\x89PNG'):
                                                mime_type = 'image/png'
                                            elif image_data.startswith(b'GIF'):
                                                mime_type = 'image/gif'
                                            
                                            image_src = f"data:{mime_type};base64,{image_base64}"
                                            print(f"✅ Extracted fill.blipFill from shape {shape_idx} ({len(image_data)} bytes, {mime_type})")
                        
                        elif hasattr(shape, 'fill') and hasattr(shape.fill, 'type') and shape.fill.type == 3:
                            if hasattr(shape.fill, 'image') and shape.fill.image:
                                image_data = shape.fill.image.blob
                                if image_data and len(image_data) > 0:
                                    image_base64 = base64.b64encode(image_data).decode('utf-8')
                                    
                                    # Determine MIME type
                                    mime_type = 'image/jpeg'
                                    if image_data.startswith(b'\x89PNG'):
                                        mime_type = 'image/png'
                                    elif image_data.startswith(b'GIF'):
                                        mime_type = 'image/gif'
                                    
                                    image_src = f"data:{mime_type};base64,{image_base64}"
                                    print(f"✅ Extracted fill image from shape {shape_idx} ({len(image_data)} bytes, {mime_type})")
                        
                        elif hasattr(shape, 'image_part') and shape.image_part:
                            image_data = shape.image_part.blob
                            if image_data and len(image_data) > 0:
                                image_base64 = base64.b64encode(image_data).decode('utf-8')
                                
                                # Determine MIME type
                                mime_type = 'image/jpeg'
                                if image_data.startswith(b'\x89PNG'):
                                    mime_type = 'image/png'
                                elif image_data.startswith(b'GIF'):
                                    mime_type = 'image/gif'
                                
                                image_src = f"data:{mime_type};base64,{image_base64}"
                                print(f"✅ Extracted image_part from shape {shape_idx} ({len(image_data)} bytes)")
                            
                    except Exception as e:
                        print(f"⚠️ Could not extract image directly: {e}")
                        
                        # Fallback: try to find image from our extracted images
                        slide_images = [img_id for img_id in images.keys() if f"image_{slide_num}_" in img_id]
                        
                        if slide_images and image_index < len(slide_images):
                            actual_image_id = slide_images[image_index]
                            image_src = images[actual_image_id]
                            image_index += 1
                            print(f"✅ Assigned image {actual_image_id} to element")
                        else:
                            # Try any available image
                            if images:
                                image_src = list(images.values())[0]
                                print(f"✅ Using fallback image")
                            else:
                                image_src = ""
                                print(f"⚠️ No images available")
                    
                    # Only create image element if we have a valid src
                    if image_src and len(image_src) > 100:  # Ensure it's a valid data URL
                        element = {
                            "id": f"image-{slide_num}-{shape_idx}",
                            "type": "image",
                            "x": int(shape.left.inches * 96),
                            "y": int(shape.top.inches * 96),
                            "width": max(int(shape.width.inches * 96), 300),
                            "height": max(int(shape.height.inches * 96), 200),
                            "src": image_src,
                            "alt": f"Image from slide {slide_num + 1}",
                            "rotation": 0,
                            "zIndex": 1,
                            "selected": False
                        }
                        slide_data["elements"].append(element)
                        print(f"✅ Created image element with src length: {len(image_src)}")
                        print(f"🔍 Image src preview: {image_src[:100]}...")
                        print(f"🔍 Full image element: {element}")
                    else:
                        print(f"⚠️ Skipped image element - no valid src found")
                
                # Handle other shape types (lines, rectangles, etc.)
                else:
                    # This is a non-text, non-image shape
                    shape_type = type(shape).__name__
                    print(f"🔍 Processing {shape_type} shape {shape_idx}")
                    
                    # Extract position and size
                    x = int(shape.left.inches * 96) if hasattr(shape, 'left') and shape.left else 0
                    y = int(shape.top.inches * 96) if hasattr(shape, 'top') and shape.top else 0
                    width = int(shape.width.inches * 96) if hasattr(shape, 'width') and shape.width else 100
                    height = int(shape.height.inches * 96) if hasattr(shape, 'height') and shape.height else 50
                    
                    # Extract fill color
                    fill_color = "transparent"  # Default transparent for shapes without fill
                    try:
                        if hasattr(shape, 'fill') and shape.fill:
                            # Check if it's a solid fill
                            if hasattr(shape.fill, 'type') and shape.fill.type == 1:  # Solid fill
                                if hasattr(shape.fill, 'solid_color') and shape.fill.solid_color:
                                    color = shape.fill.solid_color
                                    if hasattr(color, 'rgb') and color.rgb:
                                        rgb = color.rgb
                                        if hasattr(rgb, 'red') and hasattr(rgb, 'green') and hasattr(rgb, 'blue'):
                                            fill_color = f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
                                elif hasattr(shape.fill, 'fore_color') and shape.fill.fore_color:
                                    color = shape.fill.fore_color
                                    if hasattr(color, 'rgb') and color.rgb:
                                        rgb = color.rgb
                                        if hasattr(rgb, 'red') and hasattr(rgb, 'green') and hasattr(rgb, 'blue'):
                                            fill_color = f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
                            else:
                                # For non-solid fills, use transparent
                                fill_color = "transparent"
                    except Exception as e:
                        print(f"⚠️ Could not extract fill color: {e}")
                        fill_color = "transparent"
                    
                    # Extract stroke color and width
                    stroke_color = "#000000"  # Default black
                    stroke_width = 1
                    try:
                        if hasattr(shape, 'line') and shape.line:
                            if hasattr(shape.line, 'color') and shape.line.color:
                                if hasattr(shape.line.color, 'rgb') and shape.line.color.rgb:
                                    rgb = shape.line.color.rgb
                                    if hasattr(rgb, 'red') and hasattr(rgb, 'green') and hasattr(rgb, 'blue'):
                                        stroke_color = f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
                            if hasattr(shape.line, 'width') and shape.line.width:
                                stroke_width = int(shape.line.width.pt) if hasattr(shape.line.width, 'pt') else 1
                    except Exception as e:
                        print(f"⚠️ Could not extract stroke properties: {e}")
                    
                    # Determine element type based on shape
                    element_type = "rectangle"  # Default
                    if shape_type == "Line" or (hasattr(shape, 'shape_type') and shape.shape_type == 1):
                        element_type = "line"
                    elif shape_type == "Rectangle" or (hasattr(shape, 'shape_type') and shape.shape_type == 1):
                        element_type = "rectangle"
                    elif shape_type == "Oval" or (hasattr(shape, 'shape_type') and shape.shape_type == 9):
                        element_type = "circle"
                    
                    # Create the element
                    element = {
                        "id": f"{element_type}-{slide_num}-{shape_idx}",
                        "type": element_type,
                        "x": x,
                        "y": y,
                        "width": width,
                        "height": height,
                        "fill": fill_color,
                        "stroke": stroke_color,
                        "strokeWidth": stroke_width,
                        "rotation": 0,
                        "zIndex": 1,
                        "selected": False
                    }
                    
                    slide_data["elements"].append(element)
                    print(f"✅ Created {element_type} element: pos({x},{y}) size({width}x{height}) fill({fill_color}) stroke({stroke_color})")
            
            # Add any unmapped images to this slide
            for img in images:
                if img.get('slide_index') == slide_num:
                    # Check if this image is already added as an element
                    image_already_added = any(
                        element.get('type') == 'image' and 
                        element.get('x') == img['x'] and 
                        element.get('y') == img['y']
                        for element in slide_data['elements']
                    )
                    
                    if not image_already_added:
                        element = {
                            "id": f"image-{slide_num}-{img.get('shape_index', 'unmapped')}",
                            "type": "image",
                            "x": img['x'],
                            "y": img['y'],
                            "width": img['width'],
                            "height": img['height'],
                            "src": img['src'],
                            "alt": f"Image from slide {slide_num + 1}",
                            "rotation": 0,
                            "zIndex": 1,
                            "selected": False
                        }
                        slide_data["elements"].append(element)
                        print(f"✅ Added unmapped image to slide {slide_num}: pos({img['x']},{img['y']}) size({img['width']}x{img['height']})")
            
            slides.append(slide_data)
        
        # Get presentation dimensions
        slide_width = int(prs.slide_width.inches * 96) if hasattr(prs, 'slide_width') and prs.slide_width else 960
        slide_height = int(prs.slide_height.inches * 96) if hasattr(prs, 'slide_height') and prs.slide_height else 540
        
        result = {
            "title": os.path.basename(pptx_path).replace('.pptx', ''),
            "slides": slides,
            "metadata": {
                "total_slides": len(slides),
                "total_images": len(images),
                "slide_width": slide_width,
                "slide_height": slide_height,
                "presentation_width": slide_width,
                "presentation_height": slide_height
            }
        }
        
        # Optimize content for Firebase size limits
        optimized_result, final_size_mb = optimize_content_for_firebase(result, max_size_mb=0.9)
        
        return optimized_result
        
    except Exception as e:
        print(f"❌ Error parsing PPTX: {e}")
        return {
            "title": "Error",
            "slides": [],
            "metadata": {"error": str(e)}
        }

@app.route('/api/parse-pptx', methods=['POST'])
def parse_pptx():
    """Parse uploaded PPTX file"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not file.filename.lower().endswith('.pptx'):
            return jsonify({'error': 'File must be a PPTX file'}), 400
        
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            file.save(tmp_file.name)
            tmp_path = tmp_file.name
        
        try:
            # Parse the PPTX file
            result = parse_pptx_to_json(tmp_path)
            return jsonify(result)
        finally:
            # Clean up temporary file
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
            
    except Exception as e:
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@app.route('/api/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({'status': 'healthy', 'message': 'PPTX Parser API is running'})

@app.route('/api/debug-shapes', methods=['POST'])
def debug_shapes():
    """Debug endpoint to show what shapes are in a PPTX file"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file uploaded'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not file.filename.lower().endswith('.pptx'):
            return jsonify({'error': 'File must be a PPTX file'}), 400
        
        # Save uploaded file temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            file.save(tmp_file.name)
            tmp_path = tmp_file.name
        
        try:
            prs = Presentation(tmp_path)
            debug_info = {
                'total_slides': len(prs.slides),
                'slides': []
            }
            
            for slide_num, slide in enumerate(prs.slides):
                slide_info = {
                    'slide_number': slide_num + 1,
                    'total_shapes': len(slide.shapes),
                    'shapes': []
                }
                
                for shape_idx, shape in enumerate(slide.shapes):
                    shape_info = {
                        'index': shape_idx,
                        'type': type(shape).__name__,
                        'has_image': hasattr(shape, 'image') and shape.image is not None,
                        'has_fill': hasattr(shape, 'fill'),
                        'fill_type': getattr(shape.fill, 'type', None) if hasattr(shape, 'fill') else None,
                        'has_image_part': hasattr(shape, 'image_part') and shape.image_part is not None,
                        'has_text': hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text.strip(),
                        'text_content': shape.text_frame.text.strip()[:50] if hasattr(shape, 'text_frame') and shape.text_frame and shape.text_frame.text.strip() else None
                    }
                    slide_info['shapes'].append(shape_info)
                
                debug_info['slides'].append(slide_info)
            
            return jsonify(debug_info)
            
        finally:
            # Clean up temporary file
            if os.path.exists(tmp_path):
                os.unlink(tmp_path)
            
    except Exception as e:
        return jsonify({'error': f'Server error: {str(e)}'}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5001)